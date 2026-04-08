"""Calibris — Yönetici Sunum Slaytı (Neden Kullanmalıyız?)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Renk Paleti ──────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0E, 0x11, 0x17)
BG_CARD   = RGBColor(0x16, 0x1C, 0x27)
ACCENT    = RGBColor(0x4F, 0x8E, 0xF7)
ACCENT_LT = RGBColor(0x6D, 0xA8, 0xFF)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
WHITE     = RGBColor(0xE8, 0xEA, 0xF0)
GRAY      = RGBColor(0x7E, 0x8F, 0xA4)
LIGHT     = RGBColor(0xC8, 0xCD, 0xD8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Yardımcı Fonksiyonlar ────────────────────────────────────────────────────
def _dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def _tb(slide, left, top, w, h, text, sz=18, color=WHITE,
        bold=False, align=PP_ALIGN.LEFT, name="Segoe UI"):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tf


def _bullets(slide, left, top, w, h, items, sz=15, color=LIGHT,
             spacing=Pt(10)):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing


def _card(slide, left, top, w, h, fill=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_line(slide, left, top, w):
    _card(slide, left, top, w, Inches(0.04), ACCENT)


def _icon_card(slide, left, top, w, h, icon, title, desc,
               icon_clr=ACCENT):
    _card(slide, left, top, w, h)
    _tb(slide, left + Inches(0.3), top + Inches(0.25),
        Inches(0.6), Inches(0.5), icon, sz=26, color=icon_clr, bold=True)
    _tb(slide, left + Inches(0.3), top + Inches(0.7),
        w - Inches(0.5), Inches(0.4), title, sz=15, color=WHITE, bold=True)
    _tb(slide, left + Inches(0.3), top + Inches(1.1),
        w - Inches(0.5), h - Inches(1.3), desc, sz=11, color=GRAY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 1 — KAPAK
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_accent_line(s, Inches(4.5), Inches(1.8), Inches(4.3))
_tb(s, Inches(2), Inches(2.2), Inches(9.3), Inches(1),
    "Calibris", sz=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
_tb(s, Inches(2), Inches(3.3), Inches(9.3), Inches(0.7),
    "Kredi Riski Modellemede Uçtan Uca Analiz Platformu",
    sz=20, color=ACCENT_LT, align=PP_ALIGN.CENTER)
_accent_line(s, Inches(4.5), Inches(4.3), Inches(4.3))
_tb(s, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
    "Geliştirme  ·  İzleme  ·  Raporlama",
    sz=16, color=GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 2 — MEVCUT ZORLUKLAR
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_tb(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
    "Mevcut Süreçteki Zorluklar", sz=32, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3))

problems = [
    ("⏱", "Tekrarlayan İş Yükü",
     "Her yeni model çalışmasında veri çekme, temizleme ve "
     "değişken analizi adımları sıfırdan tekrarlanıyor."),
    ("🔀", "Standart Eksikliği",
     "Farklı analistler farklı yöntem ve araçlar kullanıyor; "
     "sonuçları karşılaştırmak güçleşiyor."),
    ("📋", "Manuel İzleme",
     "Model performansı Excel tablolarıyla takip ediliyor; "
     "güncel olmayan raporlar risk oluşturuyor."),
    ("🤝", "Bilgi Siloları",
     "Model konfigürasyonları kişisel bilgisayarlarda kalıyor; "
     "kişi değiştiğinde bilgi kaybı yaşanıyor."),
]
for i, (icon, title, desc) in enumerate(problems):
    col, row = i % 2, i // 2
    _icon_card(s,
               Inches(0.8) + col * Inches(6.1),
               Inches(1.8) + row * Inches(2.6),
               Inches(5.6), Inches(2.2),
               icon, title, desc, icon_clr=ORANGE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 3 — CALIBRIS NE SUNUYOR?
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_tb(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
    "Calibris Ne Sunuyor?", sz=32, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3))
_tb(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
    "Veri keşfinden model izlemeye kadar tüm süreci tek platformda birleştirir.",
    sz=15, color=GRAY)

features = [
    ("📊", "Otomatik Veri Profili",
     "Veri yüklenir yüklenmez eksik değer analizi, "
     "dağılım, IV skoru ve korelasyon hazır."),
    ("🎯", "Akıllı Değişken Seçimi",
     "Gini, Zero-Gain, RFE ve performans eğrisi ile "
     "gereksiz değişkenler otomatik elenir."),
    ("⚙️", "Esnek Modelleme",
     "LightGBM, XGBoost, Lojistik Regresyon — "
     "parametreler elle veya Optuna ile optimize edilir."),
    ("📈", "Canlı İzleme Paneli",
     "PSI, Gini/KS trendi, temerrüt oranı, "
     "göç matrisi ve backtesting tek ekranda."),
    ("💾", "Profil Kaydet & Paylaş",
     "Tüm konfigürasyon kaydedilir; ekip "
     "arkadaşı aynı profili yükleyip devam eder."),
    ("🔒", "Standart Metodoloji",
     "Herkes aynı metrikleri, aynı kuralları kullanır "
     "— denetim ve karşılaştırma kolaylaşır."),
]
for i, (icon, title, desc) in enumerate(features):
    col, row = i % 3, i // 3
    _icon_card(s,
               Inches(0.5) + col * Inches(4.2),
               Inches(2.3) + row * Inches(2.4),
               Inches(3.8), Inches(2.0),
               icon, title, desc, icon_clr=GREEN)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 4 — GELİŞTİRME AKIŞI (5 adım)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_tb(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
    "Model Geliştirme Akışı", sz=32, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3))

steps = [
    ("1", "Veri Yükle",   "SQL veya CSV ile\nçoklu tablo birleştir"),
    ("2", "Keşfet",       "Dağılım, IV, korelasyon\nve outlier analizi"),
    ("3", "Değişken Ele", "Otomatik pipeline ile\nen uygun setini seç"),
    ("4", "Modelle",      "Algoritma seç,\nOptuna ile optimize et"),
    ("5", "Kaydet",       "Profil olarak sakla\nve ekiple paylaş"),
]

card_w = Inches(2.2)
gap = Inches(0.35)
start = Inches(0.5)
top = Inches(2.2)

for i, (num, title, desc) in enumerate(steps):
    left = start + i * (card_w + gap)
    _card(s, left, top, card_w, Inches(3.0))
    _tb(s, left + Inches(0.8), top + Inches(0.3),
        Inches(0.6), Inches(0.6), num,
        sz=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    _tb(s, left + Inches(0.15), top + Inches(1.0),
        card_w - Inches(0.3), Inches(0.5), title,
        sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(s, left + Inches(0.15), top + Inches(1.6),
        card_w - Inches(0.3), Inches(1.2), desc,
        sz=12, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        _tb(s, left + card_w + Inches(0.05), top + Inches(1.2),
            Inches(0.3), Inches(0.4), "→",
            sz=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

_tb(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
    "Her adım isteğe bağlı tekrarlanabilir — analist kendi hızında ilerler.",
    sz=14, color=GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 5 — İZLEME & PROFİL YÖNETİMİ
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_tb(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
    "İzleme & Profil Yönetimi", sz=32, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3))

# Sol panel — İzleme
_card(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0))
_tb(s, Inches(0.9), Inches(2.0), Inches(5), Inches(0.5),
    "📈  Canlı Model İzleme", sz=20, color=GREEN, bold=True)
_bullets(s, Inches(0.9), Inches(2.7), Inches(5), Inches(3.5), [
    "PSI ile değişken bazlı dağılım kayması tespiti",
    "Gini / KS trendi — model ayrım gücü takibi",
    "Temerrüt oranı periyodik izleme",
    "Göç matrisi ile rating geçiş analizi",
    "Backtesting — modelin gerçek performansı",
    "Referans vs. güncel veri otomatik karşılaştırma",
], sz=14)

# Sağ panel — Profil
_card(s, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0))
_tb(s, Inches(7.2), Inches(2.0), Inches(5), Inches(0.5),
    "💾  Profil Kaydet & Paylaş", sz=20, color=GREEN, bold=True)
_bullets(s, Inches(7.2), Inches(2.7), Inches(5), Inches(3.5), [
    "Tüm analiz ayarları tek profilde saklanır",
    "Segment bazlı ayrı profiller oluşturulabilir",
    "Ekip arkadaşı profili yükleyip devam eder",
    "Denetim için analiz geçmişi korunur",
    "Model bilgisi kişiye bağımlı olmaktan çıkar",
    "Standart iş akışı tüm ekipte ortak hale gelir",
], sz=14)


# ═════════════════════════════════════════════════════════════════════════════
#  SLAYT 6 — KAPANIŞ → DEMO
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(s)

_accent_line(s, Inches(4.5), Inches(1.8), Inches(4.3))
_tb(s, Inches(2), Inches(2.2), Inches(9.3), Inches(0.8),
    "Kazanımlar", sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

_bullets(s, Inches(3.5), Inches(3.2), Inches(6.3), Inches(2.8), [
    "Tekrarlayan manuel işlemler ortadan kalkar",
    "Tüm ekip aynı standart metodoloji ile çalışır",
    "Model izleme otomatikleşir, riskler erken tespit edilir",
    "Bilgi transferi profil paylaşımı ile kolaylaşır",
    "Denetim ve raporlama için hazır altyapı sağlar",
], sz=17, color=LIGHT, spacing=Pt(14))

_accent_line(s, Inches(4.5), Inches(6.0), Inches(4.3))
_tb(s, Inches(2), Inches(6.3), Inches(9.3), Inches(0.5),
    "Şimdi canlı demo ile birlikte görelim  →",
    sz=20, color=ACCENT_LT, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
out = "outputs/Calibris_Sunum.pptx"
prs.save(out)
print(f"Sunum olusturuldu: {out}")
